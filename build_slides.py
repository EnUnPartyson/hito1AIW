"""
Generates IIT414W_DemoDay_ArielDavid.pptx — 5 content slides + 1 backup.
Run with: conda run -n iit414w python build_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK       = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
BLUE       = RGBColor(0x16, 0x48, 0x9E)   # top10 color
ORANGE     = RGBColor(0xE2, 0x6B, 0x0A)   # top5 color
ACCENT     = RGBColor(0xE8, 0x40, 0x00)   # red-orange for highlights
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0x76, 0x76, 0x76)
GREEN      = RGBColor(0x37, 0x86, 0x44)
RED_WARN   = RGBColor(0xC0, 0x00, 0x00)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def header_bar(slide, title, subtitle=None):
    """Dark header bar at top of slide."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.3), fill_color=DARK)
    add_text(slide, title,
             Inches(0.35), Inches(0.12), Inches(10), Inches(0.75),
             font_size=Pt(30), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.82), Inches(11), Inches(0.42),
                 font_size=Pt(15), color=RGBColor(0xCC, 0xCC, 0xCC), italic=True)


def footer(slide, text="IIT414W · F1 Race Strategy Advisor · Ariel Van Kilsdonk & David Hernández"):
    add_text(slide, text,
             Inches(0.3), Inches(7.15), Inches(12.5), Inches(0.3),
             font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.CENTER)


def slide_number(slide, n):
    add_text(slide, str(n),
             Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.35),
             font_size=Pt(11), color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Context + Decision
# ─────────────────────────────────────────────────────────────────────────────
def slide1(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)

    # Left navy panel
    add_rect(s, Inches(0), Inches(0), Inches(5.2), H, fill_color=DARK)

    # Race flag accent strip
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill_color=ACCENT)

    # F1 domain label
    add_text(s, "F1 RACE STRATEGY ADVISOR",
             Inches(0.35), Inches(0.5), Inches(4.6), Inches(0.5),
             font_size=Pt(11), bold=True,
             color=RGBColor(0xE8, 0x8A, 0x00), align=PP_ALIGN.LEFT)

    # Main title
    add_text(s, "Which pit strategy\nmaximizes finishing\noutcome?",
             Inches(0.35), Inches(1.1), Inches(4.5), Inches(2.4),
             font_size=Pt(32), bold=True, color=WHITE)

    # User / time / decision bullets
    bullets = [
        ("👤  User",         "Race engineer · pit wall"),
        ("⏱  Decision window", "Pre-race briefing + laps 20–30"),
        ("🎯  Decision",     "1-stop vs 2-stop tyre strategy"),
    ]
    y = Inches(3.7)
    for label, value in bullets:
        add_text(s, label,
                 Inches(0.35), y, Inches(1.8), Inches(0.4),
                 font_size=Pt(12), bold=True,
                 color=RGBColor(0xAA, 0xCC, 0xFF))
        add_text(s, value,
                 Inches(2.15), y, Inches(2.8), Inches(0.4),
                 font_size=Pt(12), color=WHITE)
        y += Inches(0.52)

    # Verdict box — right panel
    add_rect(s, Inches(5.4), Inches(1.1), Inches(7.6), Inches(2.4),
             fill_color=RGBColor(0xF8, 0xF8, 0xF8),
             line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(1))

    add_text(s, "THE FINDING",
             Inches(5.7), Inches(1.25), Inches(6.8), Inches(0.4),
             font_size=Pt(11), bold=True, color=ACCENT)

    add_text(s,
             "In wet street-circuit conditions, the strategy\n"
             "that maximizes P(top 10) is NOT the same as\n"
             "the strategy that maximizes P(top 5).\n\n"
             "A team optimizing only for points survival\n"
             "may be giving up top-5 upside.",
             Inches(5.7), Inches(1.65), Inches(6.8), Inches(1.8),
             font_size=Pt(15.5), color=DARK)

    # Decision sentence box
    add_rect(s, Inches(5.4), Inches(3.8), Inches(7.6), Inches(1.1),
             fill_color=DARK)
    add_text(s,
             '"The decision this tool supports is: which stop count gives the\n'
             'best probability of finishing top 10 vs top 5 — and do they agree?"',
             Inches(5.6), Inches(3.9), Inches(7.2), Inches(0.9),
             font_size=Pt(13), italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Two targets
    add_rect(s, Inches(5.4), Inches(5.2), Inches(3.6), Inches(1.7),
             fill_color=RGBColor(0xE8, 0xF0, 0xFB),
             line_color=BLUE, line_width=Pt(1.5))
    add_text(s, "TARGET 1", Inches(5.6), Inches(5.3), Inches(3.2), Inches(0.38),
             font_size=Pt(10), bold=True, color=BLUE)
    add_text(s, "is_top10\nPoints survival (positions 1–10)",
             Inches(5.6), Inches(5.65), Inches(3.2), Inches(0.9),
             font_size=Pt(13), color=DARK)

    add_rect(s, Inches(9.3), Inches(5.2), Inches(3.7), Inches(1.7),
             fill_color=RGBColor(0xFD, 0xF0, 0xE4),
             line_color=ORANGE, line_width=Pt(1.5))
    add_text(s, "TARGET 2", Inches(9.5), Inches(5.3), Inches(3.3), Inches(0.38),
             font_size=Pt(10), bold=True, color=ORANGE)
    add_text(s, "is_top5\nTop-end conversion (positions 1–5)",
             Inches(9.5), Inches(5.65), Inches(3.3), Inches(0.9),
             font_size=Pt(13), color=DARK)

    footer(s)
    slide_number(s, 1)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Approach
# ─────────────────────────────────────────────────────────────────────────────
def slide2(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    header_bar(s, "Approach", "Data · Split · Model · Calibration")

    # ── Data box ──
    add_rect(s, Inches(0.3), Inches(1.5), Inches(3.9), Inches(2.1),
             fill_color=LIGHT_GRAY,
             line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.75))
    add_text(s, "📊  DATA", Inches(0.5), Inches(1.6), Inches(3.5), Inches(0.4),
             font_size=Pt(12), bold=True, color=DARK)
    add_text(s,
             "2,447 driver-race entries\n"
             "6 seasons · 2019–2024\n"
             "47 raw columns",
             Inches(0.5), Inches(2.0), Inches(3.5), Inches(1.4),
             font_size=Pt(14), color=DARK)

    # ── Split table ──
    add_rect(s, Inches(4.4), Inches(1.5), Inches(4.5), Inches(2.1),
             fill_color=LIGHT_GRAY,
             line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.75))
    add_text(s, "📅  TEMPORAL SPLIT (LOCKED)", Inches(4.6), Inches(1.6),
             Inches(4.1), Inches(0.4), font_size=Pt(12), bold=True, color=DARK)
    rows = [("TRAIN",  "2019–2021", "1,132 rows", BLUE),
            ("CALIB",  "2022",      "426 rows",   MID_GRAY),
            ("TEST",   "2023–2024", "889 rows",   GREEN)]
    ry = Inches(2.05)
    for label, yrs, n, col in rows:
        add_rect(s, Inches(4.55), ry, Inches(0.85), Inches(0.42), fill_color=col)
        add_text(s, label, Inches(4.57), ry+Inches(0.05), Inches(0.82), Inches(0.35),
                 font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, yrs,  Inches(5.45), ry+Inches(0.05), Inches(1.3), Inches(0.35),
                 font_size=Pt(12), color=DARK)
        add_text(s, n,    Inches(6.8),  ry+Inches(0.05), Inches(1.9), Inches(0.35),
                 font_size=Pt(12), color=MID_GRAY)
        ry += Inches(0.48)

    # ── Model box ──
    add_rect(s, Inches(9.1), Inches(1.5), Inches(3.9), Inches(2.1),
             fill_color=LIGHT_GRAY,
             line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.75))
    add_text(s, "🤖  MODEL", Inches(9.3), Inches(1.6), Inches(3.5), Inches(0.4),
             font_size=Pt(12), bold=True, color=DARK)
    add_text(s,
             "Calibrated Logistic Regression\n"
             "One-hot categorical features\n"
             "Platt scaling on 2022 block only",
             Inches(9.3), Inches(2.0), Inches(3.5), Inches(1.4),
             font_size=Pt(14), color=DARK)

    # ── Leakage / feature row ──
    add_rect(s, Inches(0.3), Inches(3.8), Inches(12.7), Inches(0.45),
             fill_color=DARK)
    add_text(s, "FEATURE CLASSIFICATION",
             Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.36),
             font_size=Pt(11), bold=True, color=WHITE)

    feat_boxes = [
        ("PRE-RACE PREDICTORS",  "qualifying_position · constructor_tier · circuit_type",       BLUE),
        ("SCENARIO INPUTS\n(what-if controls)", "n_stops · compound_sequence",                  ORANGE),
        ("AUDIT-ONLY\n(never model inputs)",    "weather_actual · safety_car · pace aggregates", MID_GRAY),
    ]
    fx = Inches(0.3)
    for title, content, col in feat_boxes:
        add_rect(s, fx, Inches(4.35), Inches(4.1), Inches(1.7),
                 fill_color=WHITE,
                 line_color=col, line_width=Pt(2))
        add_text(s, title, fx+Inches(0.15), Inches(4.42),
                 Inches(3.8), Inches(0.55),
                 font_size=Pt(10), bold=True, color=col)
        add_text(s, content, fx+Inches(0.15), Inches(4.95),
                 Inches(3.8), Inches(0.9),
                 font_size=Pt(12), color=DARK)
        fx += Inches(4.3)

    # ── Calibration note ──
    add_rect(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.85),
             fill_color=RGBColor(0xE8, 0xF0, 0xFB),
             line_color=BLUE, line_width=Pt(1))
    add_text(s,
             "Calibration: model fit on 2019–2021 → Platt sigmoid on 2022 → test set 2023–2024 "
             "evaluated ONCE.  All reported probabilities are calibrated.",
             Inches(0.5), Inches(6.28), Inches(12.3), Inches(0.65),
             font_size=Pt(13), color=DARK)

    footer(s); slide_number(s, 2)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Results
# ─────────────────────────────────────────────────────────────────────────────
def slide3(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    header_bar(s, "Results", "Calibrated model vs baselines · test set 2023–2024")

    # ── Metrics table ──
    col_headers = ["Target", "Model", "Brier ↓", "ROC-AUC ↑", "vs Baseline"]
    col_w = [Inches(1.5), Inches(2.5), Inches(1.2), Inches(1.3), Inches(2.1)]
    col_x = [Inches(0.3)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    # Header row
    ry = Inches(1.5)
    rh = Inches(0.42)
    for i, (hdr, cx, cw) in enumerate(zip(col_headers, col_x, col_w)):
        add_rect(s, cx, ry, cw-Inches(0.04), rh, fill_color=DARK)
        add_text(s, hdr, cx+Inches(0.06), ry+Inches(0.06), cw-Inches(0.12), rh-Inches(0.1),
                 font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Data rows
    data_rows = [
        ("is_top10",  "Calibrated\nLogistic Reg.", "0.1447", "0.8726",
         "Beats heuristic\nDoes NOT beat docent floor (0.1320)",
         BLUE, RED_WARN),
        ("is_top5",   "Calibrated\nLogistic Reg.", "0.0958", "0.9217",
         "Beats heuristic baseline (0.1227)\n▲ −0.0269 Brier",
         ORANGE, GREEN),
    ]
    ry += rh
    for target, model, brier, auc, vs, tcol, vcol in data_rows:
        row_h = Inches(0.75)
        add_rect(s, col_x[0], ry, col_w[0]-Inches(0.04), row_h,
                 fill_color=RGBColor(0xE8,0xF0,0xFB) if tcol==BLUE else RGBColor(0xFD,0xF0,0xE4),
                 line_color=tcol, line_width=Pt(1.5))
        add_text(s, target, col_x[0]+Inches(0.06), ry+Inches(0.15),
                 col_w[0]-Inches(0.12), Inches(0.45),
                 font_size=Pt(13), bold=True, color=tcol, align=PP_ALIGN.CENTER)
        for val, cx, cw in zip([model, brier, auc, vs],
                                col_x[1:], col_w[1:]):
            fc = vcol if val == vs else WHITE
            add_rect(s, cx, ry, cw-Inches(0.04), row_h,
                     fill_color=LIGHT_GRAY,
                     line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
            add_text(s, val, cx+Inches(0.06), ry+Inches(0.08),
                     cw-Inches(0.12), row_h-Inches(0.1),
                     font_size=Pt(11), color=fc if val==vs else DARK,
                     align=PP_ALIGN.CENTER, bold=(val==vs))
        ry += row_h + Inches(0.08)

    # Docent floor note
    add_rect(s, Inches(0.3), ry+Inches(0.05), Inches(8.5), Inches(0.45),
             fill_color=RGBColor(0xFF,0xEB,0xEB), line_color=RED_WARN, line_width=Pt(1))
    add_text(s, "⚠  is_top10 does NOT beat the docent floor (0.1320). We report this honestly.",
             Inches(0.5), ry+Inches(0.1), Inches(8.1), Inches(0.35),
             font_size=Pt(12), bold=True, color=RED_WARN)

    # Calibration plot — embed image
    try:
        s.shapes.add_picture("calibration_curve_top10.png",
                              Inches(9.05), Inches(1.45), Inches(4.0), Inches(3.3))
        add_text(s, "Fig. 1 — is_top10 reliability diagram",
                 Inches(9.05), Inches(4.72), Inches(4.0), Inches(0.3),
                 font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)
    except Exception:
        add_text(s, "[calibration_curve_top10.png]",
                 Inches(9.05), Inches(2.5), Inches(4.0), Inches(1.0),
                 font_size=Pt(12), color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Key interpretation
    add_rect(s, Inches(0.3), Inches(5.5), Inches(8.5), Inches(1.65),
             fill_color=WHITE,
             line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(1))
    add_text(s, "WHAT THIS MEANS FOR A STRATEGY ENGINEER",
             Inches(0.5), Inches(5.58), Inches(8.1), Inches(0.38),
             font_size=Pt(11), bold=True, color=DARK)
    add_text(s,
             "• is_top10 — competitive but below docent floor: use for scenario comparison only, not deployment.\n"
             "• is_top5 — meaningfully beats the heuristic: the expansion target adds real signal beyond the grid-rule.",
             Inches(0.5), Inches(5.95), Inches(8.1), Inches(1.1),
             font_size=Pt(13), color=DARK)

    footer(s); slide_number(s, 3)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — The Trade-off
# ─────────────────────────────────────────────────────────────────────────────
def slide4(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    header_bar(s, "The Strategy Trade-off",
               "Wet street circuit · P1 grid · Top constructor — where the two targets disagree")

    # Context strip
    add_rect(s, Inches(0.3), Inches(1.55), Inches(12.7), Inches(0.5),
             fill_color=RGBColor(0xF0, 0xF4, 0xFF))
    add_text(s,
             "Scenario: qualifying P1 · top constructor · street circuit · wet race   "
             "(RF stress-test; weather_actual used as conditioning slice, not model feature)",
             Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4),
             font_size=Pt(12), color=MID_GRAY, italic=True)

    # Strategy probability table
    strategies = ["no_stop", "one_stop", "two_stop", "three_plus_stop"]
    p10 = [0.2499, 0.7196, 0.7251, 0.7150]
    p5  = [0.2722, 0.6014, 0.5956, 0.5714]

    col_w_s = [Inches(2.1), Inches(2.2), Inches(2.2), Inches(0.65), Inches(2.2), Inches(2.2)]
    # headers
    hdrs = ["Strategy", "P(top10)", "P(top5)", "", "P(top10)", "P(top5)"]

    # Table header
    tx = Inches(0.3)
    ty = Inches(2.2)
    row_h = Inches(0.52)

    hdr_labels = ["Strategy", "P(top 10)  ·  is_top10", "P(top 5)  ·  is_top5"]
    hdr_x =      [Inches(0.3), Inches(2.6), Inches(7.55)]
    hdr_w =      [Inches(2.2), Inches(4.7), Inches(4.7)]
    hdr_c =      [DARK,        BLUE,        ORANGE]
    for lbl, hx, hw, hc in zip(hdr_labels, hdr_x, hdr_w, hdr_c):
        add_rect(s, hx, ty, hw-Inches(0.05), row_h, fill_color=hc)
        add_text(s, lbl, hx+Inches(0.1), ty+Inches(0.1), hw-Inches(0.2), row_h-Inches(0.1),
                 font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    ty += row_h
    best10_idx = p10.index(max(p10))
    best5_idx  = p5.index(max(p5))

    for i, (strat, v10, v5) in enumerate(zip(strategies, p10, p5)):
        is_best10 = (i == best10_idx)
        is_best5  = (i == best5_idx)
        row_bg = LIGHT_GRAY if i % 2 == 0 else WHITE

        add_rect(s, Inches(0.3), ty, Inches(2.15), row_h, fill_color=row_bg,
                 line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
        add_text(s, strat, Inches(0.45), ty+Inches(0.1), Inches(1.95), row_h-Inches(0.12),
                 font_size=Pt(13), bold=True, color=DARK)

        # P(top10) bar area
        bar_max_w = Inches(4.4)
        bx = Inches(2.55)
        bar_bg = RGBColor(0xE0,0xEA,0xF9) if is_best10 else LIGHT_GRAY
        add_rect(s, bx, ty, Inches(4.7)-Inches(0.05), row_h, fill_color=bar_bg,
                 line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
        bar_w = bar_max_w * v10
        add_rect(s, bx+Inches(0.1), ty+Inches(0.1), bar_w, row_h-Inches(0.22),
                 fill_color=BLUE if not is_best10 else RGBColor(0x0D,0x2D,0x7A))
        label = f"  {v10:.4f}"
        if is_best10: label += "  ★ BEST"
        add_text(s, label, bx+Inches(0.15), ty+Inches(0.12),
                 Inches(4.4), row_h-Inches(0.2),
                 font_size=Pt(12), bold=is_best10,
                 color=WHITE if v10 > 0.3 else DARK)

        # P(top5) bar area
        bx2 = Inches(7.5)
        bar_bg2 = RGBColor(0xFD,0xED,0xDC) if is_best5 else LIGHT_GRAY
        add_rect(s, bx2, ty, Inches(4.7)-Inches(0.05), row_h, fill_color=bar_bg2,
                 line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
        bar_w2 = bar_max_w * v5
        add_rect(s, bx2+Inches(0.1), ty+Inches(0.1), bar_w2, row_h-Inches(0.22),
                 fill_color=ORANGE if not is_best5 else RGBColor(0xA0,0x40,0x00))
        label2 = f"  {v5:.4f}"
        if is_best5: label2 += "  ★ BEST"
        add_text(s, label2, bx2+Inches(0.15), ty+Inches(0.12),
                 Inches(4.4), row_h-Inches(0.2),
                 font_size=Pt(12), bold=is_best5,
                 color=WHITE if v5 > 0.3 else DARK)

        ty += row_h

    # Disagreement callout
    add_rect(s, Inches(0.3), ty+Inches(0.25), Inches(12.7), Inches(1.45),
             fill_color=DARK)
    add_text(s,
             "⚡  DISAGREEMENT:   "
             "is_top10 prefers two_stop (0.7251)  ·  "
             "is_top5 prefers one_stop (0.6014)\n"
             "A team optimizing only for points survival would choose the wrong strategy for top-5 upside.\n"
             "The heuristic grid-rule cannot expose this — it assigns the same probability regardless of stop count.",
             Inches(0.5), ty+Inches(0.32), Inches(12.3), Inches(1.2),
             font_size=Pt(13.5), color=WHITE, bold=False)

    footer(s); slide_number(s, 4)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Verdict + Honesty
# ─────────────────────────────────────────────────────────────────────────────
def slide5(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK)

    # Accent strip
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill_color=ACCENT)

    # Title
    add_text(s, "Verdict & Deployment Conditions",
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.65),
             font_size=Pt(28), bold=True, color=WHITE)

    # Recommendation box
    add_rect(s, Inches(0.3), Inches(1.05), Inches(12.7), Inches(1.55),
             fill_color=RGBColor(0x0D,0x2D,0x50),
             line_color=BLUE, line_width=Pt(2))
    add_text(s, "RECOMMENDATION",
             Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.38),
             font_size=Pt(11), bold=True, color=BLUE)
    add_text(s,
             "Use this tool for structured scenario comparison when you need to distinguish between "
             "the points-survival-optimal strategy and the top-5-optimal strategy.\n"
             "Do not use as a single-number oracle. Use alongside pace data and expert judgment.",
             Inches(0.5), Inches(1.48), Inches(12.3), Inches(1.0),
             font_size=Pt(15), color=WHITE)

    # Honesty sentence
    add_rect(s, Inches(0.3), Inches(2.8), Inches(12.7), Inches(0.42),
             fill_color=ACCENT)
    add_text(s, "WE DO NOT RECOMMEND DEPLOYING THIS TOOL UNLESS:",
             Inches(0.5), Inches(2.86), Inches(12.3), Inches(0.34),
             font_size=Pt(13), bold=True, color=WHITE)

    conditions = [
        ("1", "Re-evaluated on ≥1 future season beyond 2024 achieving Brier ≤ 0.1320 on is_top10"),
        ("2", "Scenario sensitivity validated on ≥3 real race disagreement cases where our preferred "
               "strategy differed from the observed team decision"),
        ("3", "Calibration for wet and midfield slices remains within ±0.02 Brier of the overall "
               "test-set score on the new evaluation season"),
    ]
    cy = Inches(3.35)
    for num, text in conditions:
        add_rect(s, Inches(0.3), cy, Inches(0.55), Inches(0.72), fill_color=ACCENT)
        add_text(s, num, Inches(0.3), cy+Inches(0.12), Inches(0.55), Inches(0.5),
                 font_size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(0.9), cy, Inches(12.1), Inches(0.72),
                 fill_color=RGBColor(0x22, 0x22, 0x44),
                 line_color=RGBColor(0x55, 0x55, 0x77), line_width=Pt(0.5))
        add_text(s, text, Inches(1.05), cy+Inches(0.1), Inches(11.8), Inches(0.6),
                 font_size=Pt(13.5), color=WHITE)
        cy += Inches(0.82)

    # Limitations strip
    add_rect(s, Inches(0.3), cy+Inches(0.1), Inches(12.7), Inches(0.48),
             fill_color=RGBColor(0x11, 0x11, 0x22))
    add_text(s,
             "Key limitations: strategy confounding with car pace  ·  "
             "no qualifying time gaps (column empty)  ·  "
             "safety car binary only  ·  2019–2024 era only",
             Inches(0.5), cy+Inches(0.17), Inches(12.3), Inches(0.38),
             font_size=Pt(11.5), color=MID_GRAY, italic=True)

    footer(s); slide_number(s, 5)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Backup (Q&A support)
# ─────────────────────────────────────────────────────────────────────────────
def slide6(prs):
    s = blank_slide(prs)
    fill_bg(s, WHITE)
    header_bar(s, "Backup — Slice Analysis & Calibration is_top5",
               "Show only if asked · not part of the 7-minute pitch")

    # Brier slice tables
    slices = {
        "Strategy type": [
            ("no_stop",         "0.0442", "0.0281"),
            ("one_stop",        "0.1559", "0.1045"),
            ("two_stop",        "0.1350", "0.0899"),
            ("three_plus_stop", "0.1521", "0.0967"),
        ],
        "Constructor tier": [
            ("backmarker", "0.1295", "0.0078"),
            ("front",      "0.1121", "0.1972"),
            ("midfield",   "0.1700", "0.1210"),
        ],
        "Weather": [
            ("dry", "0.1451", "0.0923"),
            ("wet", "0.1425", "0.1183"),
        ],
    }

    sx = Inches(0.3)
    for title, rows in slices.items():
        add_text(s, title, sx, Inches(1.55), Inches(3.8), Inches(0.38),
                 font_size=Pt(12), bold=True, color=DARK)
        # mini table header
        ty = Inches(1.92)
        for hdr, hcol, hw in [("Slice", DARK, Inches(1.6)),
                               ("top10", BLUE, Inches(1.1)),
                               ("top5", ORANGE, Inches(1.1))]:
            add_rect(s, sx if hdr=="Slice" else sx+Inches(1.65) if hdr=="top10" else sx+Inches(2.8),
                     ty, hw-Inches(0.04), Inches(0.35), fill_color=hcol)
            add_text(s, hdr,
                     sx if hdr=="Slice" else sx+Inches(1.65) if hdr=="top10" else sx+Inches(2.8),
                     ty+Inches(0.04), hw-Inches(0.08), Inches(0.28),
                     font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        ty += Inches(0.35)
        for i, (slc, b10, b5) in enumerate(rows):
            bg = LIGHT_GRAY if i%2==0 else WHITE
            is_worst_10 = (b10 == max(r[1] for r in rows))
            is_worst_5  = (b5  == max(r[2] for r in rows))
            add_rect(s, sx, ty, Inches(1.6)-Inches(0.04), Inches(0.4), fill_color=bg,
                     line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
            add_text(s, slc, sx+Inches(0.05), ty+Inches(0.06), Inches(1.5), Inches(0.3),
                     font_size=Pt(11), color=DARK)
            for val, vx, is_worst in [(b10, sx+Inches(1.65), is_worst_10),
                                       (b5,  sx+Inches(2.80), is_worst_5)]:
                vc = RED_WARN if is_worst else DARK
                add_rect(s, vx, ty, Inches(1.1)-Inches(0.04), Inches(0.4),
                         fill_color=RGBColor(0xFF,0xEE,0xEE) if is_worst else bg,
                         line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
                add_text(s, val, vx+Inches(0.05), ty+Inches(0.06), Inches(1.0), Inches(0.3),
                         font_size=Pt(11), color=vc, bold=is_worst, align=PP_ALIGN.CENTER)
            ty += Inches(0.4)
        sx += Inches(4.35)

    # is_top5 calibration plot
    try:
        s.shapes.add_picture("calibration_curve_top5.png",
                              Inches(0.3), Inches(4.1), Inches(4.5), Inches(3.0))
        add_text(s, "Fig. 2 — is_top5 reliability diagram (test 2023–2024)",
                 Inches(0.3), Inches(7.05), Inches(4.5), Inches(0.28),
                 font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)
    except Exception:
        pass

    # Worst-slice note
    add_rect(s, Inches(5.0), Inches(4.1), Inches(8.0), Inches(1.8),
             fill_color=DARK)
    add_text(s, "WORST SLICES  (Q&A ready)",
             Inches(5.2), Inches(4.18), Inches(7.6), Inches(0.38),
             font_size=Pt(12), bold=True, color=ACCENT)
    add_text(s,
             "• is_top10 hardest: midfield (Brier 0.1700) — model near decision boundary\n"
             "• is_top5 hardest: front constructors (Brier 0.1972) — overconfidence near podium\n"
             "• Both targets degrade in wet conditions — strategy confounding is strongest there",
             Inches(5.2), Inches(4.58), Inches(7.6), Inches(1.2),
             font_size=Pt(13), color=WHITE)

    footer(s); slide_number(s, 6)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# BUILD
# ─────────────────────────────────────────────────────────────────────────────
prs = new_prs()
slide1(prs)
slide2(prs)
slide3(prs)
slide4(prs)
slide5(prs)
slide6(prs)

out = "IIT414W_DemoDay_ArielDavid.pptx"
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
